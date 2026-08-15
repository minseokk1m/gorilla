#!/usr/bin/env python3
# raw 이미지 + cfg 대사 -> 편집가능한 텍스트박스가 깔린 PPTX (Canva import용)
import json, os
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

BASE = os.path.expanduser("~/Desktop/고릴라헌터스_Ep1_식자")
RAW  = os.path.join(BASE, "raw")

FONTS = {
    "B":  "/Users/minseokkim/Library/Fonts/NanumSquareOTF_acB.otf",
    "EB": "/Users/minseokkim/Library/Fonts/NanumSquareOTF_acEB.otf",
    "R":  "/Users/minseokkim/Library/Fonts/NanumSquareOTF_acR.otf",
    "L":  "/Users/minseokkim/Library/Fonts/NanumSquareOTF_acL.otf",
    "HG": "/Users/minseokkim/Library/Fonts/NanumBarunGothicBold.ttf",
}
FAMILY = {"B":"NanumSquare_ac","EB":"NanumSquare_ac","R":"NanumSquare_ac","L":"NanumSquare_ac","HG":"NanumBarunGothic"}
BOLD   = {"B":True,"EB":True,"R":False,"L":False,"HG":True}

PAGES = [
    ("p00_title.png","cfg_p00_title.json"),
    ("p01_02.png","cfg_p01_02.json"),
    ("p03_04.png","cfg_p03_04.json"),
    ("p05_06.png","cfg_p05_06.json"),
    ("p07_08.png","cfg_p07_08.json"),
    ("p09_10.png","cfg_p09_10.json"),
    ("p11.png","cfg_p11.json"),
    ("p12.png","cfg_p12.json"),
    ("p13.png","cfg_p13.json"),
    ("p14.png","cfg_p14.json"),
]

EMU_PER_PX = 9525  # 96 dpi
def PX(v): return Emu(int(round(v*EMU_PER_PX)))

# 슬라이드 기준 캔버스 = 821 x 1916 (다수 페이지 native)
CANVAS_W, CANVAS_H = 821, 1916

def wrap(text, font, max_w):
    """letter.py 와 동일한 줄바꿈 로직 (PIL 측정)."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB",(10,10)); d = ImageDraw.Draw(img)
    out = []
    for para in text.split("\n"):
        if para == "":
            out.append(""); continue
        words = para.split(" "); cur=""
        for w in words:
            trial = w if cur=="" else cur+" "+w
            if d.textlength(trial,font=font) <= max_w: cur=trial
            else:
                if cur: out.append(cur)
                if d.textlength(w,font=font) <= max_w: cur=w
                else:
                    piece=""
                    for ch in w:
                        if d.textlength(piece+ch,font=font)<=max_w: piece+=ch
                        else: out.append(piece); piece=ch
                    cur=piece
        out.append(cur)
    return out

prs = Presentation()
prs.slide_width  = PX(CANVAS_W)
prs.slide_height = PX(CANVAS_H)
blank = prs.slide_layouts[6]

from PIL import Image as PILImage
for img_name, cfg_name in PAGES:
    img_path = os.path.join(RAW, img_name)
    iw, ih = PILImage.open(img_path).size
    # fit-to-width 스케일 + 세로 중앙 정렬
    scale = CANVAS_W / iw
    disp_h = ih * scale
    off_x = 0
    off_y = (CANVAS_H - disp_h) / 2.0

    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img_path, PX(off_x), PX(off_y), PX(CANVAS_W), PX(disp_h))

    with open(os.path.join(BASE, cfg_name), encoding="utf-8") as f:
        items = json.load(f)

    for it in items:
        size = it.get("size",30)
        fkey = it.get("font","B")
        max_w = it.get("w",200)
        color = it.get("color","#141414").lstrip("#")
        rot   = it.get("rot",0)
        pilf  = ImageFont.truetype(FONTS.get(fkey,FONTS["B"]), size)
        lines = wrap(it["text"], pilf, max_w)
        # 측정
        from PIL import Image as _I, ImageDraw as _D
        _d = _D.Draw(_I.new("RGB",(10,10)))
        line_w = max((_d.textlength(l,font=pilf) for l in lines), default=10)
        lh = 1.18
        block_h = len(lines)*size*lh
        box_w = line_w*1.35 + size*1.2   # 대체폰트 대비 여유폭
        box_h = block_h + size*0.4

        # cfg 좌표/크기를 슬라이드 좌표로 변환
        cx = off_x + it["cx"]*scale
        cy = off_y + it["cy"]*scale
        bw = box_w*scale
        bh = box_h*scale
        left = cx - bw/2
        top  = cy - bh/2
        fsize = size*scale*0.75  # px -> pt

        tb = slide.shapes.add_textbox(PX(left), PX(top), PX(bw), PX(bh))
        tf = tb.text_frame
        tf.word_wrap = False   # 의도한 줄바꿈 유지, 대체폰트 재줄바꿈 방지
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in (tf.margin_left,):
            pass
        tf.margin_left = Pt(1); tf.margin_right = Pt(1)
        tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = line
            run.font.size = Pt(max(fsize,1))
            run.font.bold = BOLD.get(fkey,True)
            run.font.name = FAMILY.get(fkey,"NanumSquare_ac")
            run.font.color.rgb = RGBColor.from_string(color.upper())
        if rot:
            tb.rotation = -rot  # PIL(ccw+) -> pptx(cw+)

out = os.path.join(BASE, "gorilla_ep1_canva.pptx")
prs.save(out)
print("saved", out)
print("slides", len(prs.slides._sldIdLst))
