#!/usr/bin/env python3
# raw_v3 35장 + 대사 -> 편집가능 텍스트박스 PPTX (Canva v3 import용)
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

BASE = os.path.expanduser("~/Desktop/고릴라헌터스_Ep1_식자")
RAW  = os.path.join(BASE, "raw_v3")
EMU = 9525
def PX(v): return Emu(int(round(v*EMU)))
W, H = 1024, 1536

# 역할별 스타일
ROLE = {
 'title':    dict(size=40, font='EB', color='F5EFE2', w=860),
 'subtitle': dict(size=22, font='B',  color='E3BD55', w=620),
 'narr':     dict(size=18, font='B',  color='1A1A1A', w=600),
 'speech':   dict(size=20, font='B',  color='1A1A1A', w=400),
 'think':    dict(size=19, font='B',  color='1A1A1A', w=360),
 'caption':  dict(size=15, font='B',  color='333333', w=440),
 'sfx':      dict(size=30, font='EB', color='222222', w=260),
 'emph':     dict(size=24, font='EB', color='E3BD55', w=640),
 'label':    dict(size=16, font='B',  color='222222', w=320),
}
FAMILY = {'EB':'NanumSquare_ac','B':'NanumSquare_ac'}

# (panel, hint L/C/R, role, text)  — 한 튜플 = 한 말풍선/박스
DLG = {
 'p01':[(1,'C','title','EPISODE 01\n내가 할 수 있는 단 하나의 투자법\n고릴라 게임이란 무엇인가\n시즌 1 · 정글에 입장하다'),
        (2,'L','narr','2026년, 1년 만에 +28%\n거래는 단 5번, 매일 차트도 안 봤다'),
        (2,'C','speech','1년 전엔\n반토막이었는데…'),
        (2,'C','caption','대체 그 사이에 뭐가 바뀌었을까')],
 'p02':[(1,'C','narr','나배움, 39세'),
        (2,'C','narr','평범한 회사원,\n특별할 것 없는 하루'),
        (3,'C','narr','남들 다 한다는 주식,\n나도 작년에 처음 시작했다')],
 'p03':[(1,'C','think','이 정도면 괜찮겠지'),
        (2,'L','narr','몇 달 만에, 모은 돈이 반토막'),
        (2,'C','think','…이게 맞나?')],
 'p04':[(1,'C','narr','지푸라기라도 잡고 싶었다'),
        (1,'C','narr','주식 좀 한다는\n대학 동기 셋을 불러냈다'),
        (2,'C','narr','셋 다 자기 방식에\n확신이 넘치는 친구들')],
 'p05':[(1,'L','caption','가치파 친구 · 싸게 사서 기다린다'),
        (1,'C','speech','주식은 말이야,\nPER 낮은 걸 싸게 사서\n기다리는 거야'),
        (2,'C','speech','지금 AI주 중에\n저평가된 거, 그거 사')],
 'p06':[(1,'L','caption','성장파 친구 · 잘 나가는 데 올라탄다'),
        (1,'C','speech','무슨 소리야,\n매출 40% 느는 데를 사야지'),
        (2,'L','speech','싼 데는 싼 이유가\n있는 거라고'),
        (2,'R','speech','비싸게 사면 물려')],
 'p07':[(1,'C','caption','한탕수 · 핫한 것만 쫓는 후배'),
        (1,'C','speech','에이 형들,\n지금 핫한 게 답이에요!'),
        (2,'C','speech','요즘 다들\n이거 산다니까요?'),
        (3,'C','think','다… 다 맞는 말 같은데?')],
 'p08':[(1,'C','narr','셋 다 맞는 말 같았다'),
        (2,'C','speech','그럼… 셋 다\n사면 되잖아?'),
        (2,'C','caption','그게 첫 번째 실수였다')],
 'p09':[(1,'C','sfx','탭 탭 탭'),
        (2,'C','think','이제 됐다'),
        (3,'C','narr','그리고, 며칠 뒤')],
 'p10':[(1,'C','sfx','우르르… 쾅'),
        (1,'C','speech','셋 다… 반토막?'),
        (2,'C','think','대체 왜 나만\n당하는 거야…?')],
 'p11':[(2,'C','emph','진짜 이유\n프로는 같은 데이터로 너보다 빠르다'),
        (2,'L','narr','같은 게임으론,\n일반 직장인은 절대 못 이긴다'),
        (1,'C','think','난 처음부터\n질 수밖에 없었던 건가…')],
 'p12':[(1,'L','narr','그날 퇴근길,\n발걸음이 유난히 무거웠다'),
        (2,'C','narr','홀린 듯,\n불 켜진 서점으로 들어섰다')],
 'p13':[(1,'C','think','고릴라… 게임?'),
        (2,'C','think','투자랑 고릴라가\n무슨 상관이지')],
 'p14':[(1,'C','narr','책 뒤표지에 작은 안내가 있었다\n「고릴라 헌터스」, 매주 모이는 투자 스터디'),
        (2,'C','think','이런 데를… 내가?')],
 'p15':[(1,'C','narr','며칠을 망설이다,\n나는 그 문을 두드렸다'),
        (2,'C','narr','여기서부터, 모든 게 달라졌다')],
 'p16':[(1,'C','speech','어, 새 분 오셨네'),
        (2,'L','caption','좌장 구본질 · 반도체 엔지니어 출신 투자자'),
        (2,'C','speech','『고릴라 게임』 들고 왔네,\n잘 찾아왔어'),
        (3,'C','speech','여기… 주식 공부하는\n모임이라고 해서요')],
 'p17':[(1,'C','speech','좋다는 거 다 샀는데,\n다 같이 반토막이 났어요'),
        (2,'C','speech','근데, 자네가\n왜 졌는지는 아나?'),
        (2,'C','speech','…그걸 알면 안 졌겠죠')],
 'p18':[(1,'C','speech','자네가 본 건 전부 ‘결과’야,\nPER·매출·주가'),
        (1,'C','speech','다들 이 결과만 보고 사고팔지'),
        (2,'C','speech','우린 ‘구조’를 봐\n왜 이 회사가 구조적으로 안 무너지는가'),
        (3,'C','think','결과가 아니라… 구조?')],
 'p19':[(1,'C','speech','결과 싸움은 자네가\n프로를 절대 못 이겨'),
        (2,'C','speech','근데 구조는 달라, 느려도 돼\n한번 정해진 판은 한동안 안 바뀌니까'),
        (2,'C','think','그래서 내가 졌던 거구나')],
 'p20':[(1,'C','speech','세상 사람을 산업 지식과\n주식 지식, 두 축으로 나눠 보자'),
        (2,'C','think','난 둘 다 어중간한데…\n어디지?')],
 'p21':[(1,'C','speech','VC도, 펀드매니저도,\n완전 초보도 아닌\n산업도 주식도 ‘중간’인 사람\n그 정중앙이 고릴라 게임 투자자야'),
        (1,'L','emph','일반인이 프로보다 잘할 수 있는,\n거의 유일한 자리'),
        (2,'C','speech','그게… 저 같은\n사람이라고요?')],
 'p22':[(1,'C','speech','제가… 그 자리에\n있다고요?'),
        (2,'C','speech','오! 그럼 저도\n고릴라 투자자네요?'),
        (2,'R','speech','넌 ‘배울 자세’부터,\n핫한 거 말고'),
        (3,'C','speech','아 그건 좀…')],
 'p23':[(1,'C','speech','정중앙의 조건은 딱 둘이야\n하나, 배울 자세. 둘, 꾸준함'),
        (2,'C','speech','머리가 좋아야 하는 게 아니야,\n핫한 걸 참을 줄 알아야 해'),
        (1,'R','think','그거라면…\n나도 할 수 있지 않을까')],
 'p24':[(1,'C','speech','한번 표준이 되면\n한동안 안 무너져\n사방이 이 회사 위에서 돌거든,\n그게 고릴라야'),
        (2,'C','think','사방이 다\nNVIDIA 위에서 돈다고…?')],
 'p25':[(1,'C','narr','정글에도 왕이 있다,\n모두가 그 밑에서 산다'),
        (2,'C','think','난 왕이 아니라,\n그냥 핫한 동물들을 산 거였구나')],
 'p26':[(1,'C','speech','그러니까 고릴라 게임은…'),
        (2,'C','caption','고릴라 = 표준 잡은 1등에 집중\n결과가 아니라 ‘구조’를 본다')],
 'p27':[(1,'C','speech','근데 유튜브엔 매일\n차트 분석 채널이 엄청 많던데\n그게 더 정석 아닌가요?'),
        (2,'C','speech','그게 바로 ‘결과’를\n매일 쳐다보는 게임이야\n우린 그거 안 해')],
 'p28':[(1,'L','label','데이트레이딩·모멘텀·밸류\n(매일 차트, 결과를 쫓는다)'),
        (1,'R','label','고릴라 게임\n(구조가 정해지는 순간만 읽는다)'),
        (2,'C','narr','이건 데이트레이딩도, 모멘텀도, 밸류도 아니야\n한번 사면 토네이도가 끝날 때까지 길게 가는 게임이지')],
 'p29':[(1,'C','caption','결정이 적다 (연 1~2번) · 장기 보유 (직장과 양립) · 프로가 못 하는 게임'),
        (1,'C','speech','결정은 적게, 보유는 길게\n프로는 매분기 성과를 내야 해서\n이걸 못 해')],
 'p30':[(1,'C','emph','일반인이 프로보다 잘할 수 있는,\n거의 유일한 영역'),
        (1,'C','speech','프로가 못 하는 이 자리,\n그게 자네 자리야'),
        (2,'C','narr','처음으로, 내가 이길 수 있는\n판이 보였다')],
 'p31':[(1,'C','speech','에이 그래도 핫한 게 최고죠\n2배 레버리지 들어갑니다!'),
        (2,'L','speech','그렇게 빚내서 쫓다가,\n매번 물리는 거야'),
        (2,'R','think','또 시작이네…')],
 'p32':[(1,'C','speech','레버리지는 토네이도가\n끝난 자리에서 사람을 삼켜'),
        (2,'C','speech','에이, 저는 괜찮아요'),
        (2,'C','caption','그 말을, 한탕수는 곧 후회하게 된다')],
 'p33':[(1,'C','think','가치, 성장, 핫한 거…\n난 결과만 쫓았네'),
        (2,'C','think','구조라…')],
 'p34':[(1,'C','speech','내가 산 건…\n고릴라가 아니었구나'),
        (2,'C','think','이제, 진짜를 사고 싶다')],
 'p35':[(1,'C','emph','이 만화의 주인공은 너야'),
        (1,'L','speech','검증된 것만 사고,\n매일 차트 볼 시간은 없고\n그래도 자산을 불리고 싶은 사람'),
        (3,'C','speech','이제 정글의 법칙을 하나씩 배우자\n근데 이 법칙, 사실 70년 전\n아이오와 옥수수밭에서 시작됐어'),
        (3,'C','caption','Episode 1 끝 · 다음 화 · 기술은 어떻게 퍼지는가\n(아이오와 옥수수밭에서 ChatGPT까지)')],
}

# 설명 캡션 띠 (하단 파치먼트 박스) — 핵심 개념 장
BANDS = {
 'p16':'「고릴라 헌터스」, 결과가 아니라 ‘구조’를 보는 투자 스터디\n매주 모여 한 기업이 왜 안 무너지는지를 함께 뜯어본다',
 'p18':'결과 = PER·매출·주가처럼 눈에 보이는 숫자\n구조 = 왜 이 기업이 쉽게 안 무너지는가',
 'p20':'가로축 = 주식 지식, 세로축 = 산업 지식\n9칸으로 투자자를 나눠 본다',
 'p21':'둘 다 ‘중간’인 정중앙, 거기가 프로도 초보도 아닌 일반인의 자리',
 'p24':'표준(de facto standard) = 모두가 그 위에서 일하게 된 1등\n한번 표준이 되면 쉽게 안 바뀐다, 그 1등이 ‘고릴라’',
 'p28':'매일 차트 보는 게임이 아니라\n산업 구조가 바뀌는 드문 순간에만 길게 베팅하는 게임',
 'p29':'① 결정이 적다 (연 1~2번)  ② 길게 보유 (직장과 양립)  ③ 프로는 단기 성과 압박 때문에 못 한다',
}

HX = {'L':300, 'C':512, 'R':724}

def layout(items):
    total = max(2, max(it[0] for it in items))
    band = H/total
    by = {}
    for it in items: by.setdefault(it[0], []).append(it)
    out = []
    for panel, its in by.items():
        y0 = (panel-1)*band
        n = len(its)
        for i,(p,hint,role,text) in enumerate(its):
            frac = 0.18 + (0.6*(i/(n-1)) if n>1 else 0.12)
            cy = y0 + band*frac
            cx = HX.get(hint,512)
            out.append((text,cx,cy,role))
    return out

prs = Presentation()
prs.slide_width = PX(W); prs.slide_height = PX(H)
blank = prs.slide_layouts[6]

for i in range(1,36):
    pid = f"p{i:02d}"
    img = os.path.join(RAW, pid+".png")
    iw,ih = PILImage.open(img).size
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, 0, 0, PX(W), PX(H))
    for text,cx,cy,role in layout(DLG[pid]):
        st = ROLE[role]
        nlines = text.count("\n")+1
        bw = st['w']; bh = (st['size']*1.35)*nlines + st['size']
        left = cx - bw/2; top = cy - bh/2
        left = max(8, min(left, W-bw-8)); top = max(8, min(top, H-bh-8))
        tb = slide.shapes.add_textbox(PX(left), PX(top), PX(bw), PX(bh))
        tf = tb.text_frame; tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ('margin_left','margin_right','margin_top','margin_bottom'):
            setattr(tf, m, Pt(1))
        for j,line in enumerate(text.split("\n")):
            par = tf.paragraphs[0] if j==0 else tf.add_paragraph()
            par.alignment = PP_ALIGN.CENTER
            run = par.add_run(); run.text = line
            run.font.size = Pt(st['size']*0.75)
            run.font.bold = True
            run.font.name = FAMILY[st['font']]
            run.font.color.rgb = RGBColor.from_string(st['color'])
    # 하단 설명 캡션 띠
    band = BANDS.get(pid)
    if band:
        nln = band.count("\n")+1
        bw = W-100; bh = 30*nln+34
        left=(W-bw)/2; top=H-bh-34
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(left),PX(top),PX(bw),PX(bh))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string('F4ECD8')
        sh.line.color.rgb = RGBColor.from_string('8A7A55'); sh.line.width = Pt(1.2)
        sh.shadow.inherit = False
        tf = sh.text_frame; tf.word_wrap = True
        for m in ('margin_left','margin_right','margin_top','margin_bottom'):
            setattr(tf, m, Pt(6))
        for j,line in enumerate(band.split("\n")):
            par = tf.paragraphs[0] if j==0 else tf.add_paragraph()
            par.alignment = PP_ALIGN.CENTER
            run = par.add_run(); run.text = line
            run.font.size = Pt(15*0.95); run.font.bold = True
            run.font.name = 'NanumSquare_ac'
            run.font.color.rgb = RGBColor.from_string('2A2018')

out = os.path.join(BASE, "gorilla_ep1_v3.pptx")
prs.save(out)
print("saved", out, "slides", len(prs.slides._sldIdLst))
